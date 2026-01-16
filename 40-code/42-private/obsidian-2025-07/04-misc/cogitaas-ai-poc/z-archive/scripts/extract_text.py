#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx>=0.6.0",
# ]
# ///
import sys
from pptx import Presentation


def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    extracted_text = ""
    for slide_number, slide in enumerate(presentation.slides):
        extracted_text += f"\nSlide {slide_number + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
    return extracted_text


if __name__ == "__main__":
    pptx_path = sys.argv[1]
    text = extract_text_from_pptx(pptx_path)
    with open("parsed_data/CSF_Project_workflow/text.txt", "w") as f:
        f.write(text)
